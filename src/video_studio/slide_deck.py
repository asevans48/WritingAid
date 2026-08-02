"""Slide-deck editor helpers: build a project from chapter scenes,
distribute timings from a pasted script, and stitch the result
into an MP4 (image stills + per-slide audio).

Kept separate from the chapter-deck export module because the
slide editor's model is finer-grained — one slide per action
favorite — and the audio handling is per-slide rather than a
single master mix.
"""

from __future__ import annotations

import re
import subprocess
from pathlib import Path
from typing import Any, List, Optional, Tuple

from src.video_studio.models import (
    SlideDeckProject, SlideGroup, SlidePage, TitleCard,
)
from src.video_studio.stitcher import (
    ffmpeg_available, stitch_clips, stitch_with_transitions,
)


# Minimum duration we clamp a slide to. Anything shorter and the
# stitcher's image-to-MP4 conversion produces noisy output, and
# the audio mux can drop the slide entirely on some ffmpeg builds.
MIN_SLIDE_SECONDS = 1.0

# How long the LAST slide of a group may keep showing after its
# narration has finished. Kept tiny so groups butt-join cleanly on
# a cut instead of freezing on the final image for ~0.5s (the "gap
# between groups" writers see when no transition is selected). When
# the narration runs longer than the slide, the slide still
# stretches to cover it — this only caps the trailing silent hold.
MAX_TRAILING_HOLD_SECONDS = 0.15


def build_slide_deck_from_chapter(
    chapter_scenes: List[Any],
    working_dir: Path,
    chapter_id: str = "",
    chapter_label: str = "",
    default_duration_seconds: float = 4.0,
) -> SlideDeckProject:
    """Walk the chapter's scenes and assemble a SlideDeckProject.

    Per scene:
      * Slideshow mode → one page per action's favorite image.
      * Other modes → one page per scene's favorite clip when the
        favorite is an image (videos are skipped — the slide
        editor works with stills).
    """
    deck = SlideDeckProject(
        name=chapter_label or "Slide deck",
        chapter_id=chapter_id,
        working_dir=str(working_dir),
        wpm_estimate=150,
    )
    page_index = 0
    for scene in chapter_scenes:
        scene_label = scene.name or f"Scene {page_index + 1}"
        actions = getattr(scene, "actions", None) or []
        # Use per-action favorite images whenever the actions HAVE
        # images — not only when the scene is flagged "slideshow".
        # Favoriting or importing an action image doesn't flip the
        # scene's mode, so gating on mode alone made the editor
        # ignore favorited photos and report "no favorites."
        has_action_images = any(
            (getattr(a, "images", None) or []) for a in actions)
        # Each SCENE becomes a GROUP named after the scene (the card
        # in the video studio); each ACTION becomes a SLIDE placed on
        # that group's timeline. The group is created lazily on the
        # scene's first usable slide so scenes with no image don't
        # leave empty groups behind.
        scene_group = None
        running = 0.0
        if actions and (
                getattr(scene, "mode", "video") == "slideshow"
                or has_action_images):
            for action in actions:
                # A lone image is implicitly the favorite — persist
                # that so provenance is explicit downstream.
                try:
                    action.ensure_single_image_favorite()
                except Exception:
                    pass
                img = action.favorite_image()
                if img is None:
                    continue
                path_str = (img.file_path or "").strip()
                if not path_str:
                    continue
                p = Path(path_str)
                if not p.exists() or p.stat().st_size == 0:
                    continue
                dur = max(
                    MIN_SLIDE_SECONDS,
                    float(
                        action.display_seconds
                        or scene.image_display_seconds
                        or default_duration_seconds))
                if scene_group is None:
                    scene_group = SlideGroup(
                        name=scene_label, source_scene_id=scene.id)
                    deck.groups.append(scene_group)
                page = SlidePage(
                    index=page_index,
                    label=(action.name
                           or f"Action {page_index + 1}"),
                    image_path=str(p),
                    duration_seconds=dur,
                    source_scene_id=scene.id,
                    source_action_id=action.id,
                    text_overlay=getattr(img, "overlay", None),
                    group_id=scene_group.id,
                    start_time_seconds_in_group=round(running, 3),
                )
                scene_group.page_ids.append(page.id)
                running += dur
                deck.pages.append(page)
                page_index += 1
            continue
        # Non-slideshow scene — single slide from the favorite
        # clip when it's an image. Still gets its own scene group.
        clip = scene.favorite_clip()
        if clip is None or not clip.file_path:
            continue
        p = Path(clip.file_path)
        if (not p.exists()
                or p.stat().st_size == 0
                or p.suffix.lower() not in
                {".png", ".jpg", ".jpeg", ".webp", ".gif", ".bmp"}):
            continue
        scene_group = SlideGroup(
            name=scene_label, source_scene_id=scene.id)
        deck.groups.append(scene_group)
        page = SlidePage(
            index=page_index,
            label=scene_label,
            image_path=str(p),
            duration_seconds=max(
                MIN_SLIDE_SECONDS,
                float(
                    scene.image_display_seconds
                    or default_duration_seconds)),
            source_scene_id=scene.id,
            group_id=scene_group.id,
            start_time_seconds_in_group=0.0,
        )
        scene_group.page_ids.append(page.id)
        deck.pages.append(page)
        page_index += 1
    return deck


def rebuild_scene_group(
    deck: SlideDeckProject,
    scene: Any,
    working_dir: Path,
) -> Optional[SlideGroup]:
    """RESET one scene's contribution to the deck: drop the scene's
    existing group + slides and rebuild them fresh from the scene's
    current actions / favorites. Keeps the group in its original
    position when it had one. Returns the new group (or None when the
    scene now yields no slides). Other scenes' groups are untouched."""
    tmp = build_slide_deck_from_chapter(
        [scene], Path(working_dir),
        chapter_id=getattr(deck, "chapter_id", ""))
    sid = getattr(scene, "id", "")
    # Find where the old group sat so the rebuild lands in place.
    old_pos = next(
        (i for i, g in enumerate(deck.groups)
         if getattr(g, "source_scene_id", "") == sid),
        len(deck.groups))
    old_page_ids = {
        p.id for p in deck.pages
        if getattr(p, "source_scene_id", None) == sid}
    deck.pages = [
        p for p in deck.pages if p.id not in old_page_ids]
    deck.groups = [
        g for g in deck.groups
        if getattr(g, "source_scene_id", "") != sid]
    deck.pages.extend(tmp.pages)
    for off, g in enumerate(tmp.groups):
        deck.groups.insert(min(old_pos + off, len(deck.groups)), g)
    return tmp.groups[0] if tmp.groups else None


def sync_group_to_scene(deck: SlideDeckProject, scene: Any) -> int:
    """SYNC a scene's group edits back onto the scene: push each
    slide's duration onto its source action's ``display_seconds`` and
    the group's (possibly renamed) name back onto the scene. Returns
    the count of fields changed."""
    sid = getattr(scene, "id", "")
    actions_by_id = {
        a.id: a for a in (getattr(scene, "actions", None) or [])}
    changed = 0
    for p in deck.pages:
        if getattr(p, "source_scene_id", None) != sid:
            continue
        aid = getattr(p, "source_action_id", None)
        a = actions_by_id.get(aid) if aid else None
        if a is None:
            continue
        dur = round(float(getattr(p, "duration_seconds", 0.0) or 0.0), 2)
        if abs(float(getattr(a, "display_seconds", 0.0) or 0.0)
               - dur) > 0.01:
            a.display_seconds = dur
            changed += 1
    g = next(
        (g for g in deck.groups
         if getattr(g, "source_scene_id", "") == sid), None)
    if g is not None and g.name and g.name != getattr(
            scene, "name", ""):
        try:
            scene.name = g.name
            changed += 1
        except Exception:
            pass
    return changed


def assign_orphan_slides_to_scene_groups(
    deck: SlideDeckProject,
    scene_names_by_id: dict,
) -> int:
    """Bring a LEGACY ungrouped deck up to the scene→group /
    action→slide model: group every slide by its ``source_scene_id``
    into a group named after the scene (the video-studio card).

    Only runs when the deck has NO groups yet — a deck the writer has
    already arranged into groups is left untouched. Returns the number
    of slides grouped."""
    if getattr(deck, "groups", None):
        return 0
    by_scene: dict = {}
    grouped = 0
    for p in deck.pages:
        sid = getattr(p, "source_scene_id", None)
        if not sid:
            continue
        g = by_scene.get(sid)
        if g is None:
            name = scene_names_by_id.get(sid) or "Scene"
            g = SlideGroup(name=name, source_scene_id=sid)
            deck.groups.append(g)
            by_scene[sid] = g
        p.group_id = g.id
        if p.id not in g.page_ids:
            g.page_ids.append(p.id)
        grouped += 1
    # Place each group's slides sequentially on its timeline.
    for g in by_scene.values():
        members = [p for p in deck.pages if p.group_id == g.id]
        running = 0.0
        for p in members:
            p.start_time_seconds_in_group = round(running, 3)
            running += max(
                MIN_SLIDE_SECONDS,
                float(getattr(p, "duration_seconds", 0.0) or 0.0))
    return grouped


def suggest_timings_from_script(
    deck: SlideDeckProject,
    script_text: str,
) -> Tuple[int, str]:
    """Split the script into per-slide chunks (one blank-line
    paragraph per slide), assign each chunk to its slide, and
    compute each slide's duration from word count + WPM. Returns
    ``(slides_touched, message)``.

    Honors ``locked_duration`` — locked slides keep their
    explicit time and the script chunk still lands on them as
    text so the writer can read along.

    Pure heuristic — no LLM call. Writers can override per slide
    afterward or paste a more detailed script and re-run.
    """
    if not deck.pages:
        return (0, "No slides in deck.")
    text = (script_text or "").strip()
    if not text:
        return (0, "No script text provided.")
    # Split on blank lines. Falls back to one-chunk-per-sentence
    # when the writer pasted prose without paragraph breaks.
    chunks = [
        c.strip() for c in re.split(r"\n\s*\n", text)
        if c.strip()
    ]
    if len(chunks) <= 1:
        # Sentence split as the fallback so a single paragraph
        # still distributes across slides instead of all landing
        # on slide 1.
        sentences = re.split(r"(?<=[.!?])\s+", chunks[0] if chunks else text)
        chunks = [s.strip() for s in sentences if s.strip()]
    if not chunks:
        return (0, "No usable script chunks found.")
    pages = deck.pages
    # When the writer gave us more chunks than slides, pool the
    # tail into the last slide. When fewer chunks, distribute
    # evenly across slides.
    per_slide: List[str] = ["" for _ in pages]
    if len(chunks) <= len(pages):
        # Distribute by stretching — chunk i lands on page
        # round(i * pages / chunks). This keeps order without
        # leaving early slides starved.
        n_pages = len(pages)
        n_chunks = len(chunks)
        for i, chunk in enumerate(chunks):
            target = int(round(i * (n_pages - 1) / max(1, n_chunks - 1))) \
                if n_chunks > 1 else 0
            if per_slide[target]:
                per_slide[target] += "\n\n" + chunk
            else:
                per_slide[target] = chunk
    else:
        # More chunks than slides: assign chunk i to slide
        # floor(i * pages / chunks). Tail chunks pool on the last.
        for i, chunk in enumerate(chunks):
            target = min(
                len(pages) - 1,
                int(i * len(pages) / len(chunks)))
            if per_slide[target]:
                per_slide[target] += "\n\n" + chunk
            else:
                per_slide[target] = chunk
    wpm = max(60, int(deck.wpm_estimate or 150))
    touched = 0
    for page, chunk in zip(pages, per_slide):
        if not chunk:
            continue
        page.script_text = chunk
        if not page.locked_duration:
            words = max(1, len(chunk.split()))
            secs = max(
                MIN_SLIDE_SECONDS,
                round(words / (wpm / 60.0), 2))
            page.duration_seconds = secs
        touched += 1
    return (
        touched,
        f"Distributed {len(chunks)} chunk(s) across "
        f"{touched} slide(s) at ~{wpm} WPM.")


def adjust_slide_to_audio(page: SlidePage) -> bool:
    """When a recorded / imported audio take is attached, set the
    slide's duration to match. Honors ``locked_duration``. Returns
    True when the duration changed."""
    if page.locked_duration:
        return False
    if (page.audio_duration_seconds <= 0
            or not page.audio_path):
        return False
    new_dur = max(MIN_SLIDE_SECONDS,
                  round(page.audio_duration_seconds + 0.2, 2))
    if abs(new_dur - page.duration_seconds) < 0.05:
        return False
    page.duration_seconds = new_dur
    return True


def distribute_group_timings(
    deck: SlideDeckProject,
    group: SlideGroup,
) -> int:
    """Evenly split a group's ``target_total_seconds`` across its
    UNLOCKED pages. Locked pages keep their exact times; the
    remainder splits equally across the rest. Returns the count
    of pages whose duration changed."""
    if group.target_total_seconds <= 0:
        return 0
    pages_by_id = {p.id: p for p in deck.pages}
    members = [
        pages_by_id[pid] for pid in group.page_ids
        if pid in pages_by_id]
    if not members:
        return 0
    locked = [p for p in members if p.locked_duration]
    unlocked = [p for p in members if not p.locked_duration]
    if not unlocked:
        return 0
    locked_total = sum(p.duration_seconds for p in locked)
    remainder = max(
        0.0, group.target_total_seconds - locked_total)
    if remainder <= 0:
        return 0
    per_page = max(
        MIN_SLIDE_SECONDS, remainder / len(unlocked))
    changed = 0
    for p in unlocked:
        if abs(p.duration_seconds - per_page) > 0.05:
            p.duration_seconds = round(per_page, 2)
            changed += 1
    return changed


def export_slide_deck_to_pptx(
    deck: SlideDeckProject,
    output_path: Path,
) -> Tuple[bool, str, List[str]]:
    """Render a SlideDeckProject as a PowerPoint file.

    One slide per ``SlidePage`` with its image as the only visual.
    Per-slide audio is embedded (if present) and configured to
    play on slide entry, so when the deck is opened in PowerPoint
    / Keynote / Slides the narration runs automatically. Each
    slide's advance time is set to its ``duration_seconds`` so a
    slideshow presentation matches the recorded timings.

    No text overlays — the writer can edit freely in their slide
    tool. Returns ``(success, message, skipped)``.
    """
    try:
        from pptx import Presentation
        from pptx.util import Inches
        from pptx.dml.color import RGBColor
    except Exception as e:
        return (
            False,
            (
                "python-pptx isn't installed. Install it with:\n"
                "  pip install python-pptx\n\n"
                f"Underlying error: {e}"),
            [])
    pages = [
        p for p in deck.pages
        if p.image_path and Path(p.image_path).exists()]
    if not pages:
        return (
            False,
            "No slides with usable images on disk.",
            [])
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)
    blank_layout = prs.slide_layouts[6]
    slide_w_emu = int(prs.slide_width)
    slide_h_emu = int(prs.slide_height)
    skipped: List[str] = []
    AUDIO_EXTS = {
        ".wav", ".mp3", ".m4a", ".aac", ".ogg", ".flac",
        ".opus", ".aiff", ".aif"}
    for page in pages:
        slide = prs.slides.add_slide(blank_layout)
        # Pure black background — the image is the show.
        bg = slide.background.fill
        bg.solid()
        bg.fore_color.rgb = RGBColor(0x00, 0x00, 0x00)
        # Image, fitted preserving aspect ratio.
        try:
            left, top, width, height = _fit_image_to_slide(
                Path(page.image_path),
                slide_w_emu, slide_h_emu)
            slide.shapes.add_picture(
                page.image_path, left, top,
                width=width, height=height)
        except Exception as e:
            skipped.append(
                f"{page.label or page.id}: image embed failed "
                f"({e})")
            continue
        # Per-slide audio — embedded via add_movie. PowerPoint
        # treats audio files (.wav / .mp3 / etc.) as media objects
        # too. We tuck the speaker icon offscreen and patch the XML
        # to play automatically + hide while playing — that's
        # what writers expect for narration.
        if (page.audio_path
                and Path(page.audio_path).exists()
                and Path(page.audio_path).suffix.lower()
                in AUDIO_EXTS):
            try:
                # 32 EMU ≈ ~0.03 inch — effectively invisible.
                from pptx.util import Emu
                icon_w = Emu(304800)  # 0.5 in
                icon_h = Emu(304800)
                # Park the speaker icon in the bottom-right corner
                # off the visible canvas, so writers can still find
                # it to edit if needed.
                left = slide_w_emu - icon_w - Emu(91440)
                top = slide_h_emu - icon_h - Emu(91440)
                movie = slide.shapes.add_movie(
                    page.audio_path, left, top, icon_w, icon_h,
                    mime_type="audio/x-wav")
                # Patch the timing XML so the audio auto-plays on
                # slide entry. python-pptx leaves this as
                # ``click`` by default.
                _patch_media_autoplay(slide, movie)
            except Exception as e:
                skipped.append(
                    f"{page.label or page.id}: audio embed "
                    f"failed ({e})")
        # Per-slide advance time — slideshow honors this as the
        # auto-advance interval. The XML lives on the slide's
        # transition element.
        try:
            secs = max(
                MIN_SLIDE_SECONDS,
                float(page.duration_seconds))
            _set_slide_advance_time(slide, secs)
        except Exception as e:
            skipped.append(
                f"{page.label or page.id}: advance-time set "
                f"failed ({e})")
        # Per-slide transition effect. The first slide ignores its
        # transition (no previous slide). PowerPoint stores the
        # transition on the SLIDE you're transitioning INTO, which
        # matches the writer's "transition_in" semantic.
        if (page.index > 0
                and (page.transition_in or "cut") != "cut"):
            try:
                _set_slide_transition_effect(
                    slide,
                    page.transition_in,
                    float(page.transition_seconds or 0.7))
            except Exception as e:
                skipped.append(
                    f"{page.label or page.id}: transition set "
                    f"failed ({e})")
    output_path.parent.mkdir(parents=True, exist_ok=True)
    try:
        prs.save(str(output_path))
    except Exception as e:
        return (
            False, f"PowerPoint save failed: {e}", skipped)
    return (
        True,
        f"PowerPoint deck saved to {output_path}.",
        skipped)


def _fit_image_to_slide(
    image_path: Path, slide_w_emu: int, slide_h_emu: int,
) -> Tuple[int, int, int, int]:
    """Compute (left, top, width, height) in EMUs that fits the
    image in the slide preserving its aspect ratio. Falls back to
    full-bleed when the image can't be probed."""
    try:
        from PIL import Image
        with Image.open(image_path) as im:
            img_w, img_h = im.size
    except Exception:
        return (0, 0, slide_w_emu, slide_h_emu)
    if img_w <= 0 or img_h <= 0:
        return (0, 0, slide_w_emu, slide_h_emu)
    slide_aspect = slide_w_emu / slide_h_emu
    img_aspect = img_w / img_h
    if img_aspect > slide_aspect:
        width = slide_w_emu
        height = int(slide_w_emu / img_aspect)
        left = 0
        top = (slide_h_emu - height) // 2
    else:
        height = slide_h_emu
        width = int(slide_h_emu * img_aspect)
        top = 0
        left = (slide_w_emu - width) // 2
    return (left, top, width, height)


def _patch_media_autoplay(slide, movie) -> None:
    """Replace the click-to-play trigger on a movie shape with an
    auto-play-on-slide-entry trigger, and hide the speaker icon
    while presenting.

    PowerPoint stores this in the slide's ``<p:timing>`` element.
    python-pptx exposes the underlying XML so we can walk it,
    swap ``clickEffect`` → ``withEffect``, and add a
    ``showMediaCtrls=0`` attribute on the picture.
    """
    from pptx.oxml.ns import qn
    timing = slide.element.find(qn("p:timing"))
    if timing is None:
        return
    # Find every condition that triggers on click and flip it.
    for cond in timing.iter(qn("p:cond")):
        evt = cond.get("evt")
        if evt == "onClick":
            cond.set("evt", "afterEffect")
    # Mark every video filter so PowerPoint hides the speaker
    # icon during playback (it's already tucked offscreen, but
    # belt and braces).
    media_id = movie.shape_id
    for video_shape in slide.element.iter(qn("p:pic")):
        nvSpPr = video_shape.find(qn("p:nvPicPr"))
        if nvSpPr is None:
            continue
        cNvPr = nvSpPr.find(qn("p:cNvPr"))
        if cNvPr is None:
            continue
        if cNvPr.get("id") == str(media_id):
            # Hide media controls while the slide plays.
            cNvPr.set("hidden", "1")
            break


def _set_slide_advance_time(slide, seconds: float) -> None:
    """Configure the slide's transition so PowerPoint auto-
    advances after ``seconds`` seconds during a slideshow.

    PPT stores advance time in milliseconds on the
    ``<p:transition>`` element via the ``advTm`` attribute.
    ``advClick="0"`` keeps the click-advance off so the auto-time
    is the only trigger.
    """
    from pptx.oxml.ns import qn
    from lxml import etree
    transition = slide.element.find(qn("p:transition"))
    if transition is None:
        transition = etree.SubElement(
            slide.element, qn("p:transition"))
    transition.set("advClick", "0")
    transition.set("advTm", str(int(seconds * 1000)))


# Map our shared xfade transition keys (CHAPTER_TRANSITIONS) to
# PowerPoint transition element names. Some xfade options have
# no exact PPT analogue — we substitute the visually closest
# option so writers don't lose the cue.
_PPT_TRANSITION_MAP = {
    "cut": None,
    "fade": ("fade", {}),
    "fadeblack": ("fade", {"thruBlk": "1"}),
    "fadewhite": ("fade", {"thruBlk": "1"}),
    "dissolve": ("fade", {}),
    "slideleft": ("push", {"dir": "l"}),
    "slideright": ("push", {"dir": "r"}),
    "slideup": ("push", {"dir": "u"}),
    "slidedown": ("push", {"dir": "d"}),
    "wipeleft": ("wipe", {"dir": "l"}),
    "wiperight": ("wipe", {"dir": "r"}),
    "circleopen": ("circle", {}),
    "circleclose": ("circle", {}),
    "radial": ("wheel", {"spokes": "1"}),
}


def _set_slide_transition_effect(
    slide, xfade_key: str, seconds: float,
) -> None:
    """Attach a PowerPoint-style transition effect to the slide
    based on the writer's xfade pick. ``cut`` is a no-op.

    PowerPoint expresses transitions via a child element of
    ``<p:transition>`` (e.g. ``<p:fade/>``, ``<p:wipe dir="l"/>``).
    The transition's speed is ``"fast"``/``"med"``/``"slow"`` —
    mapped from seconds.
    """
    from pptx.oxml.ns import qn
    from lxml import etree
    mapping = _PPT_TRANSITION_MAP.get(
        (xfade_key or "cut").lower())
    if mapping is None:
        return
    transition = slide.element.find(qn("p:transition"))
    if transition is None:
        transition = etree.SubElement(
            slide.element, qn("p:transition"))
    # Pick speed based on duration.
    if seconds <= 0.4:
        speed = "fast"
    elif seconds >= 1.2:
        speed = "slow"
    else:
        speed = "med"
    transition.set("spd", speed)
    name, attrs = mapping
    effect = etree.SubElement(
        transition, qn(f"p:{name}"))
    for k, v in attrs.items():
        effect.set(k, v)


def _track_map_get(m: Any, idx: int, default: Any = None) -> Any:
    """Read a per-track dict that may be keyed by int OR str
    (JSON round-trips int keys to strings)."""
    if not m:
        return default
    if idx in m:
        return m[idx]
    if str(idx) in m:
        return m[str(idx)]
    return default


def _next_free_track(group: SlideGroup) -> int:
    """Lowest track index not already used by a clip in ``group``."""
    used = {
        int(getattr(c, "track_index", 0) or 0)
        for c in (getattr(group, "audio_clips", None) or [])}
    i = 0
    while i in used:
        i += 1
    return i


def copy_group_track(
    src_group: SlideGroup,
    track_index: int,
    dst_group: SlideGroup,
    dst_track_index: Optional[int] = None,
) -> int:
    """Replicate one lane from ``src_group`` into ``dst_group``.

    Copies every clip on ``src_group``'s ``track_index`` — with
    all its edits (gain, fades, trims, de-esser, denoise) — plus
    that lane's settings (name, gain, de-esser, mute, background
    loop). Clips are DEEP-copied with fresh ids so the two groups
    stay fully independent: editing one lane never touches the
    other. Lands on ``dst_track_index`` (or the next free lane in
    the destination when omitted). Returns the destination track
    index. Does NOT recompose — the caller owns that so it can
    batch the render.
    """
    from uuid import uuid4
    src_clips = [
        c for c in (getattr(src_group, "audio_clips", None) or [])
        if int(getattr(c, "track_index", 0) or 0) == track_index]
    if dst_track_index is None:
        dst_track_index = _next_free_track(dst_group)
    if getattr(dst_group, "audio_clips", None) is None:
        dst_group.audio_clips = []
    for c in src_clips:
        nc = c.model_copy(deep=True)
        nc.id = f"aclip_{uuid4().hex[:10]}"
        nc.track_index = int(dst_track_index)
        dst_group.audio_clips.append(nc)
    # Carry the lane's treatment across so the copy sounds the
    # same, not just plays the same clips.
    for attr in (
            "track_names", "track_gain_db",
            "track_deesser_intensity", "track_muted",
            "track_background"):
        src_map = getattr(src_group, attr, None) or {}
        val = _track_map_get(src_map, track_index, None)
        if val is None:
            continue
        dst_map = dict(getattr(dst_group, attr, None) or {})
        # Normalize to int keys so the destination stays clean.
        dst_map = {
            int(k): v for k, v in dst_map.items()
            if str(k).lstrip("-").isdigit()}
        dst_map[int(dst_track_index)] = val
        setattr(dst_group, attr, dst_map)
    return int(dst_track_index)


def group_has_background_lane(group: SlideGroup) -> bool:
    """True when ``group`` owns a background-loop lane with clips —
    i.e. a lane flagged in ``track_background`` that actually has
    audio on it. This is what the deck-wide universal bed defers to
    (OFF mode) or replaces (ON mode)."""
    bgs = getattr(group, "track_background", None) or {}
    bg_lanes = {
        int(k) for k, v in bgs.items()
        if str(k).lstrip("-").isdigit() and bool(v)}
    if not bg_lanes:
        return False
    for c in (getattr(group, "audio_clips", None) or []):
        if int(getattr(c, "track_index", 0) or 0) in bg_lanes:
            return True
    return False


def plan_universal_background_regions(
    group_timeline: List[Tuple[float, float, bool]],
    deck_dur: float,
    universal_len: float,
    complete_final: bool,
) -> Tuple[List[Tuple[float, float]], float]:
    """Decide WHERE (and for how long) the deck's background bed
    plays over the final timeline.

    ``group_timeline`` is one ``(start, end, suppress_here)`` per
    rendered segment, in final-video seconds — ``suppress_here`` is
    True where the bed must NOT play (a group with its own bed in
    gap-fill mode, or a card whose deck-background was disabled).
    The caller folds universal-vs-gap-fill INTO those flags, so this
    function just walks the maximal runs where the bed IS allowed.

    Returns ``(regions, extra_tail)`` where each region is
    ``(start_seconds, play_length_seconds)`` — the bed drops in at
    ``start`` and loops for ``play_length`` — and ``extra_tail`` is
    how many seconds the video must extend past its natural end so
    the final loop can complete (0 when cut). Loops are whole
    cycles; a mid-deck region is cut at its boundary; only the run
    that reaches the deck's end honors ``complete_final``.
    """
    import math
    regions: List[Tuple[float, float]] = []
    extra_tail = 0.0
    if universal_len <= 0 or not group_timeline or deck_dur <= 0:
        return regions, extra_tail
    runs: List[Tuple[float, float]] = []
    cur: Optional[List[float]] = None
    for (s, e, suppress_here) in group_timeline:
        if suppress_here:
            if cur is not None:
                runs.append((cur[0], cur[1]))
                cur = None
        else:
            if cur is None:
                cur = [s, e]
            else:
                cur[1] = e
    if cur is not None:
        runs.append((cur[0], cur[1]))
    for (rs, re) in runs:
        length = max(0.0, re - rs)
        if length <= 0:
            continue
        reaches_end = abs(re - deck_dur) < 0.05
        if reaches_end and complete_final:
            cycles = max(1, math.ceil(length / universal_len - 1e-6))
            full = cycles * universal_len
            extra_tail = max(extra_tail, (rs + full) - deck_dur)
            length = full
        regions.append((round(rs, 3), round(length, 3)))
    return regions, extra_tail


def compose_group_overlay_variant(
    group: SlideGroup,
    working_dir: str,
    exclude_background: bool,
) -> str:
    """Compose ``group``'s overlay, optionally EXCLUDING its
    background-loop lanes (narration / foreground only). Used when
    a universal deck bed replaces group beds: the group renders its
    voice, the deck bed carries the music. Returns the WAV path, or
    "" when there's nothing to render. Non-excluded calls just defer
    to the normal resolver so the cache stays shared."""
    if not exclude_background:
        return resolve_group_overlay(group, working_dir=working_dir)
    clips = getattr(group, "audio_clips", None) or []
    if not clips:
        return ""
    bgs = getattr(group, "track_background", None) or {}
    bg_lanes = {
        int(k) for k, v in bgs.items()
        if str(k).lstrip("-").isdigit() and bool(v)}
    if not bg_lanes:
        # No bed to strip — identical to the normal overlay.
        return resolve_group_overlay(group, working_dir=working_dir)
    try:
        from src.video_studio.audio_edit import compose_clips
    except Exception:
        return resolve_group_overlay(group, working_dir=working_dir)
    # Mute the background lanes on top of any existing mutes.
    muted = {
        int(k): bool(v)
        for k, v in (getattr(group, "track_muted", None) or {}).items()
        if str(k).lstrip("-").isdigit()}
    for lane in bg_lanes:
        muted[lane] = True
    dest_dir = Path(
        working_dir
        or (Path.home() / ".writingaid_slides")) / "group_overlay"
    try:
        dest_dir.mkdir(parents=True, exist_ok=True)
    except Exception:
        return resolve_group_overlay(group, working_dir=working_dir)
    dest = dest_dir / f"{group.id}_narration.wav"
    try:
        result = compose_clips(
            clips, dest,
            track_gain_db=getattr(group, "track_gain_db", None),
            track_deesser_intensity=getattr(
                group, "track_deesser_intensity", None),
            track_muted=muted,
            track_background=getattr(
                group, "track_background", None))
    except Exception:
        return resolve_group_overlay(group, working_dir=working_dir)
    if not result.success:
        return ""
    return str(dest)


def resolve_deck_background(deck: SlideDeckProject) -> str:
    """Return an on-disk WAV for the deck's background bed, with
    every clip edit (trim, gain, fades, de-esser, noise reduction,
    per-lane treatment) baked in. The bed is stored as a
    ``SlideGroup`` (``deck.background_group``), so this just defers
    to ``resolve_group_overlay`` — the exact renderer group tracks
    use. Returns "" when there is no bed."""
    bg_group = getattr(deck, "background_group", None)
    if bg_group is None or not (
            getattr(bg_group, "audio_clips", None) or []):
        return ""
    return resolve_group_overlay(
        bg_group, working_dir=deck.working_dir)


def resolve_group_overlay(
    group: SlideGroup,
    working_dir: str = "",
) -> str:
    """Return an on-disk overlay WAV for ``group`` that carries
    every audio-clip edit (gain, de-essing, noise reduction,
    per-lane looping, fades, trims) baked in.

    ``audio_clips`` is the source of truth; ``overlay_audio_path``
    is a rendered cache the group editor refreshes after each
    edit. Exports and previews used to trust that cache blindly —
    so when the cache file was missing (deck moved between
    machines, working_dir changed, the composed WAV was cleaned,
    or the clips were populated by a path that never opened the
    group editor) the export silently dropped the group's sound.

    This helper closes that gap: if the cached file is present it
    is returned as-is (fast path — the editor keeps it fresh);
    otherwise the overlay is recomposed from ``audio_clips`` right
    here so the narration — with all its edits — still lands in
    the render. Returns "" when the group genuinely has no audio.
    """
    cached = (getattr(group, "overlay_audio_path", "") or "").strip()
    if cached and Path(cached).exists() \
            and Path(cached).stat().st_size > 0:
        return cached
    clips = getattr(group, "audio_clips", None) or []
    if not clips:
        return ""
    # Cache missing but clips exist — recompose from the source
    # of truth so the edits aren't lost. Import lazily to keep
    # the module importable in environments without the audio
    # stack wired up.
    try:
        from src.video_studio.audio_edit import compose_clips
    except Exception:
        return cached  # best effort — nothing else we can do
    dest_dir = Path(
        working_dir
        or getattr(group, "working_dir", "")
        or (Path.home() / ".writingaid_slides")) / "group_overlay"
    try:
        dest_dir.mkdir(parents=True, exist_ok=True)
    except Exception:
        return cached
    # Deterministic name so repeated exports reuse one file
    # instead of littering the folder with timestamped renders.
    dest = dest_dir / f"{group.id}_export_overlay.wav"
    try:
        result = compose_clips(
            clips, dest,
            track_gain_db=getattr(group, "track_gain_db", None),
            track_deesser_intensity=getattr(
                group, "track_deesser_intensity", None),
            track_muted=getattr(group, "track_muted", None),
            track_background=getattr(
                group, "track_background", None))
    except Exception:
        return cached
    if not result.success:
        return cached
    # Refresh the in-memory cache so later exports in this session
    # take the fast path (and an autosave persists the valid
    # path). Guarded — some SlideGroup builds may not expose the
    # setters.
    try:
        group.overlay_audio_path = str(dest)
        group.overlay_audio_duration_seconds = float(
            result.duration_seconds or 0.0)
    except Exception:
        pass
    return str(dest)


# ---------------------------------------------------------------------
# Export compression (target file size)
# ---------------------------------------------------------------------
def estimate_deck_duration_seconds(deck: SlideDeckProject) -> float:
    """Rough total runtime of the deck from the model (no render):
    per group, the larger of its placed-slide holds or its overlay
    audio; summed across groups, minus inter-group transition
    overlaps. Good enough to size a compression target."""
    total = 0.0
    first = True
    for g in (getattr(deck, "groups", None) or []):
        placed = [
            p for p in deck.pages
            if getattr(p, "group_id", None) == g.id
            and getattr(p, "start_time_seconds_in_group", None)
            is not None]
        if not placed:
            continue
        holds = sum(
            max(0.25, float(getattr(p, "duration_seconds", 0.0)
                            or 0.0))
            for p in placed)
        overlay = float(
            getattr(g, "overlay_audio_duration_seconds", 0.0) or 0.0)
        seg = max(holds, overlay)
        if not first:
            secs = float(
                getattr(g, "inter_group_transition_seconds", 0.0)
                or 0.0)
            kind = getattr(g, "inter_group_transition_in", "cut")
            if (kind or "cut") != "cut" and secs > 0:
                seg = max(0.0, seg - secs)
        total += seg
        first = False
    # Orphan slides (no group) each add their own hold.
    group_ids = {g.id for g in (getattr(deck, "groups", None) or [])}
    for p in deck.pages:
        gid = getattr(p, "group_id", None)
        if not gid or gid not in group_ids:
            total += max(0.25, float(
                getattr(p, "duration_seconds", 0.0) or 0.0))
    return round(total, 2)


def recommend_export_target_mb(
    deck: SlideDeckProject,
    width: int = 1280,
    height: int = 720,
) -> Tuple[float, str]:
    """Heuristic compression recommendation: a target size (MB)
    that keeps a narrated slideshow crisp but shareable. Returns
    ``(target_mb, rationale)``. Deterministic — always available,
    no LLM required."""
    dur = estimate_deck_duration_seconds(deck)
    minutes = max(0.1, dur / 60.0)
    # Any video-background card pushes the per-minute budget up
    # (motion needs more bits than near-static slides).
    has_video = any(
        getattr(p, "card", None)
        and getattr(p.card, "kind", "") == "video"
        for p in deck.pages)
    px = width * height
    # MB per minute, tuned for x264 at these resolutions.
    if px >= 1920 * 1080:
        mb_per_min = 22.0 if has_video else 14.0
    elif px >= 1280 * 720:
        mb_per_min = 12.0 if has_video else 8.0
    else:
        mb_per_min = 7.0 if has_video else 5.0
    target = max(2.0, round(minutes * mb_per_min, 1))
    rationale = (
        f"~{dur:.0f}s deck at {width}×{height}"
        + (" with video background" if has_video else "")
        + f" → ~{mb_per_min:.0f} MB/min ≈ {target:.0f} MB "
          "(good balance of quality and shareable size).")
    return (target, rationale)


def compress_to_target_size(
    video_path: Path,
    target_mb: float,
    audio_kbps: int = 192,
) -> Tuple[bool, str]:
    """Re-encode ``video_path`` IN PLACE (two-pass x264) to land
    near ``target_mb`` megabytes. Splits the budget between video
    and a fixed audio bitrate; clamps the video bitrate to a sane
    floor so tiny targets still produce a watchable file."""
    if target_mb <= 0:
        return (True, "no target")
    if not ffmpeg_available():
        return (False, "ffmpeg not found on PATH.")
    dur = _probe_media_duration(video_path)
    if dur <= 0:
        return (False, "could not probe duration")
    has_audio_stream = _probe_has_audio(video_path)
    ab = audio_kbps if has_audio_stream else 0
    # Total budget → video bitrate, with a 3% container margin.
    total_kbps = (target_mb * 8.0 * 1024.0) / dur
    video_kbps = int(max(120.0, total_kbps * 0.97 - ab))
    passlog = video_path.with_name(video_path.stem + "_2pass")
    tmp = video_path.with_name(video_path.stem + "_sized.mp4")
    null_dev = "/dev/null"
    base = ["ffmpeg", "-y", "-i", str(video_path.resolve()),
            "-c:v", "libx264", "-b:v", f"{video_kbps}k",
            "-pix_fmt", "yuv420p", "-preset", "medium",
            "-passlogfile", str(passlog)]
    p1 = base + ["-pass", "1", "-an", "-f", "mp4", null_dev]
    p2 = base + ["-pass", "2"]
    if has_audio_stream:
        p2 += ["-c:a", "aac", "-b:a", f"{ab}k", "-ac", "2"]
    else:
        p2 += ["-an"]
    p2 += ["-movflags", "+faststart", str(tmp.resolve())]
    try:
        r1 = subprocess.run(
            p1, capture_output=True, text=True, timeout=1800)
        if r1.returncode != 0:
            return (False, "compress pass 1 failed:\n"
                    + (r1.stderr or "")[-300:])
        r2 = subprocess.run(
            p2, capture_output=True, text=True, timeout=1800)
        if r2.returncode != 0:
            return (False, "compress pass 2 failed:\n"
                    + (r2.stderr or "")[-300:])
    except subprocess.TimeoutExpired:
        return (False, "compression timed out.")
    finally:
        for suffix in ("-0.log", "-0.log.mbtree"):
            try:
                Path(str(passlog) + suffix).unlink(missing_ok=True)
            except Exception:
                pass
    try:
        tmp.replace(video_path)
    except Exception as e:
        return (False, f"could not finalize: {e}")
    actual_mb = video_path.stat().st_size / (1024 * 1024)
    return (True, f"compressed to {actual_mb:.1f} MB "
            f"(target {target_mb:.0f} MB, {video_kbps}k video).")


def _probe_media_duration(path: Path) -> float:
    try:
        proc = subprocess.run(
            ["ffprobe", "-v", "error",
             "-show_entries", "format=duration",
             "-of", "default=nw=1:nk=1", str(Path(path).resolve())],
            capture_output=True, text=True, timeout=30)
        return float((proc.stdout or "0").strip() or 0.0)
    except Exception:
        return 0.0


def _probe_has_audio(path: Path) -> bool:
    try:
        proc = subprocess.run(
            ["ffprobe", "-v", "error", "-select_streams", "a",
             "-show_entries", "stream=codec_type", "-of",
             "csv=p=0", str(Path(path).resolve())],
            capture_output=True, text=True, timeout=30)
        return "audio" in (proc.stdout or "")
    except Exception:
        return False


# ---------------------------------------------------------------------
# Title / ending cards
# ---------------------------------------------------------------------
_FONT_CANDIDATES = (
    "/System/Library/Fonts/Supplemental/Arial Bold.ttf",
    "/System/Library/Fonts/Supplemental/Arial.ttf",
    "/System/Library/Fonts/Helvetica.ttc",
    "/Library/Fonts/Arial.ttf",
    "/usr/share/fonts/truetype/dejavu/DejaVuSans-Bold.ttf",
    "/usr/share/fonts/truetype/dejavu/DejaVuSans.ttf",
    "/usr/share/fonts/truetype/liberation/"
    "LiberationSans-Bold.ttf",
)


def _resolve_font_file() -> str:
    """First available TTF/TTC on disk for ffmpeg ``drawtext``.
    Returns "" when none is found (caller then skips text)."""
    for cand in _FONT_CANDIDATES:
        if Path(cand).exists():
            return cand
    return ""


def _sanitize_color(value: str, default: str = "black") -> str:
    """Accept ``#RRGGBB`` / ``#RGB`` / a named color; fall back to
    ``default`` for anything ffmpeg wouldn't parse."""
    v = (value or "").strip()
    if not v:
        return default
    if re.fullmatch(r"#([0-9a-fA-F]{3}|[0-9a-fA-F]{6})", v):
        # ffmpeg wants 0xRRGGBB.
        h = v.lstrip("#")
        if len(h) == 3:
            h = "".join(c * 2 for c in h)
        return f"0x{h.upper()}"
    if re.fullmatch(r"[A-Za-z]+", v):
        return v.lower()
    return default


def _hex_to_rgba(value: str, default: str = "#FFFFFF") -> tuple:
    """Parse ``#RRGGBB`` / ``#RGB`` into an ``(r, g, b, 255)``
    tuple for PIL, falling back to ``default``."""
    v = (value or "").strip() or default
    m = re.fullmatch(r"#([0-9a-fA-F]{3}|[0-9a-fA-F]{6})", v)
    if not m:
        v = default
    h = v.lstrip("#")
    if len(h) == 3:
        h = "".join(c * 2 for c in h)
    return (int(h[0:2], 16), int(h[2:4], 16), int(h[4:6], 16), 255)


def _render_text_overlay_png(
    card: TitleCard, width: int, height: int, out_png: Path,
) -> bool:
    """Render the card's title + subtitle onto a full-frame
    TRANSPARENT PNG via PIL (this ffmpeg build has no ``drawtext``).
    The PNG is then overlaid — and faded — by ffmpeg. Returns True
    when text was drawn."""
    title = (getattr(card, "title", "") or "").strip()
    subtitle = (getattr(card, "subtitle", "") or "").strip()
    # Text elements (kind == "text"), in layer order, plus any legacy
    # text_boxes that haven't been migrated yet.
    boxes = [
        e for e in sorted(
            (getattr(card, "elements", None) or []),
            key=lambda e: getattr(e, "z", 0))
        if getattr(e, "kind", "text") == "text"
        and (getattr(e, "text", "") or "").strip()]
    boxes += [
        b for b in (getattr(card, "text_boxes", None) or [])
        if (getattr(b, "text", "") or "").strip()]
    if not (title or subtitle or boxes):
        return False
    fontfile = _resolve_font_file()
    if not fontfile:
        return False
    try:
        from PIL import Image, ImageDraw, ImageFont
    except Exception:
        return False
    scale = height / 1080.0
    title_px = max(10, int((getattr(card, "title_size", 72)
                            or 72) * scale))
    sub_px = max(8, int((getattr(card, "subtitle_size", 40)
                         or 40) * scale))
    gap = int(24 * scale)
    try:
        tfont = ImageFont.truetype(fontfile, title_px)
        sfont = ImageFont.truetype(fontfile, sub_px)
    except Exception:
        return False
    img = Image.new("RGBA", (width, height), (0, 0, 0, 0))
    draw = ImageDraw.Draw(img)

    def _measure(text: str, font) -> tuple:
        b = draw.textbbox((0, 0), text, font=font)
        return (b[2] - b[0], b[3] - b[1])

    tw, th = _measure(title, tfont) if title else (0, 0)
    sw, sh = _measure(subtitle, sfont) if subtitle else (0, 0)
    block_h = (th if title else 0) + ((sh + gap) if subtitle else 0)
    block_w = max(tw, sw)
    pos = (getattr(card, "text_position", "center")
           or "center").lower()
    if pos == "top":
        top = int(height * 0.10)
    elif pos == "bottom":
        top = int(height * 0.90 - block_h)
    else:
        top = int((height - block_h) / 2)
    # ── Effect: legibility box behind the text block ──
    box_color = (getattr(card, "text_box_color", "") or "").strip()
    if box_color:
        try:
            r, g, b, _ = _hex_to_rgba(box_color, "#000000")
            alpha = int(max(0.0, min(1.0, float(
                getattr(card, "text_box_opacity", 0.5)
                or 0.5))) * 255)
            padx = int(28 * scale)
            pady = int(18 * scale)
            bx0 = int((width - block_w) / 2) - padx
            by0 = top - pady
            bx1 = int((width + block_w) / 2) + padx
            by1 = top + block_h + pady
            draw.rectangle(
                [bx0, by0, bx1, by1], fill=(r, g, b, alpha))
        except Exception:
            pass
    outline_color = (
        getattr(card, "text_outline_color", "") or "").strip()
    outline_w = int(getattr(card, "text_outline_width", 0) or 0)
    outline_px = max(0, int(outline_w * scale))
    shadow = bool(getattr(card, "text_shadow", False))

    def _draw_line(text, font, cx_y, fill):
        w0, _h0 = _measure(text, font)
        x = (width - w0) / 2
        yy = cx_y
        # Effect: drop shadow.
        if shadow:
            off = max(2, int(3 * scale))
            draw.text(
                (x + off, yy + off), text, font=font,
                fill=(0, 0, 0, 160))
        # Effect: outline / stroke (PIL supports stroke_width).
        if outline_px > 0 and outline_color:
            draw.text(
                (x, yy), text, font=font, fill=fill,
                stroke_width=outline_px,
                stroke_fill=_hex_to_rgba(outline_color, "#000000"))
        else:
            draw.text((x, yy), text, font=font, fill=fill)

    y = top
    if title:
        _draw_line(
            title, tfont, y,
            _hex_to_rgba(
                getattr(card, "title_color", "#FFFFFF"), "#FFFFFF"))
        y += th + gap
    if subtitle:
        _draw_line(
            subtitle, sfont, y,
            _hex_to_rgba(
                getattr(card, "subtitle_color", "#DDDDDD"),
                "#DDDDDD"))
    # ── Free-placed PowerPoint-style text boxes ──
    for box in boxes:
        _draw_text_box(draw, box, width, height, scale, fontfile)
    try:
        img.save(str(out_png))
    except Exception:
        return False
    return True


def _draw_text_box(
    draw, box, width: int, height: int,
    scale: float, fontfile: str,
) -> None:
    """Draw one free-placed ``TextBox`` onto ``draw``: word-wrap the
    text to the box width, vertically align it, and apply the box's
    legibility fill / outline / shadow. Geometry is normalized (0..1)
    so it lands identically at any resolution."""
    from PIL import ImageFont
    text = (getattr(box, "text", "") or "").strip()
    if not text:
        return
    fpx = max(8, int((getattr(box, "font_size", 72) or 72) * scale))
    try:
        font = ImageFont.truetype(fontfile, fpx)
    except Exception:
        return
    bx = float(getattr(box, "x", 0.1)) * width
    by = float(getattr(box, "y", 0.1)) * height
    bw = max(1.0, float(getattr(box, "w", 0.8)) * width)
    bh = max(1.0, float(getattr(box, "h", 0.2)) * height)

    def _tw(s: str) -> float:
        b = draw.textbbox((0, 0), s, font=font)
        return b[2] - b[0]

    # Word-wrap each paragraph to the box width.
    lines: list = []
    for para in text.split("\n"):
        cur = ""
        for word in para.split(" "):
            trial = (cur + " " + word).strip()
            if not cur or _tw(trial) <= bw:
                cur = trial
            else:
                lines.append(cur)
                cur = word
        lines.append(cur)
    asc, desc = font.getmetrics()
    lh = asc + desc
    total_h = lh * max(1, len(lines))
    valign = (getattr(box, "valign", "middle") or "middle").lower()
    if valign == "top":
        ty = by
    elif valign == "bottom":
        ty = by + bh - total_h
    else:
        ty = by + (bh - total_h) / 2.0
    # Legibility fill behind the whole box.
    fill_hex = (getattr(box, "box_color", "") or "").strip()
    if fill_hex:
        try:
            r, g, b, _ = _hex_to_rgba(fill_hex, "#000000")
            alpha = int(max(0.0, min(1.0, float(
                getattr(box, "box_opacity", 0.5) or 0.5))) * 255)
            draw.rectangle(
                [bx, by, bx + bw, by + bh], fill=(r, g, b, alpha))
        except Exception:
            pass
    align = (getattr(box, "align", "center") or "center").lower()
    outline_hex = (getattr(box, "outline_color", "") or "").strip()
    outline_px = max(
        0, int((getattr(box, "outline_width", 0) or 0) * scale))
    shadow = bool(getattr(box, "shadow", False))
    fill = _hex_to_rgba(getattr(box, "color", "#FFFFFF"), "#FFFFFF")
    y = ty
    for ln in lines:
        w0 = _tw(ln)
        if align == "left":
            x = bx
        elif align == "right":
            x = bx + bw - w0
        else:
            x = bx + (bw - w0) / 2.0
        if shadow:
            off = max(2, int(3 * scale))
            draw.text(
                (x + off, y + off), ln, font=font,
                fill=(0, 0, 0, 160))
        if outline_px > 0 and outline_hex:
            draw.text(
                (x, y), ln, font=font, fill=fill,
                stroke_width=outline_px,
                stroke_fill=_hex_to_rgba(outline_hex, "#000000"))
        else:
            draw.text((x, y), ln, font=font, fill=fill)
        y += lh


def bake_text_overlay(
    image_path: str,
    overlay: TitleCard,
    out_path: Path,
) -> bool:
    """Composite ``overlay``'s styled text (with effects) directly
    onto ``image_path`` and save to ``out_path``. Used to burn a
    per-slide text overlay into the still before it flows through
    the normal image-based render. Returns True on success."""
    if overlay is None:
        return False
    try:
        from PIL import Image
    except Exception:
        return False
    try:
        base = Image.open(image_path).convert("RGBA")
    except Exception:
        return False
    w, h = base.size
    txt_png = out_path.with_name(out_path.stem + "_txt.png")
    if not _render_text_overlay_png(overlay, w, h, txt_png):
        return False
    try:
        txt = Image.open(str(txt_png)).convert("RGBA")
        combined = Image.alpha_composite(base, txt)
        combined.convert("RGB").save(str(out_path))
        ok = True
    except Exception:
        ok = False
    finally:
        try:
            txt_png.unlink(missing_ok=True)
        except Exception:
            pass
    return ok


def _video_first_frame(path: str, out_png: Path) -> bool:
    """Extract a video's first frame to ``out_png``. Best-effort."""
    try:
        subprocess.run(
            ["ffmpeg", "-y", "-i", str(Path(path).resolve()),
             "-vframes", "1", str(out_png)],
            capture_output=True, timeout=30)
        return out_png.exists()
    except Exception:
        return False


def _paste_media_element(
    base_rgba, el, width: int, height: int, out_path: Path,
) -> None:
    """Composite one image / video element onto ``base_rgba`` (a PIL
    RGBA image, mutated in place): fit the media within the element's
    normalized box, centered. A video contributes its first frame in
    the still — its motion is added by the MP4 renderer."""
    try:
        from PIL import Image, ImageOps
    except Exception:
        return
    path = getattr(el, "media_path", "") or ""
    if not path or not Path(path).exists():
        return
    tmp_frame = None
    try:
        if getattr(el, "kind", "") == "video":
            tmp_frame = out_path.with_name(
                out_path.stem + f"_el_{getattr(el, 'id', 'v')}.png")
            if not _video_first_frame(path, tmp_frame):
                return
            img = Image.open(str(tmp_frame)).convert("RGBA")
        else:
            img = Image.open(path).convert("RGBA")
    except Exception:
        return
    try:
        bx = int(float(getattr(el, "x", 0.0)) * width)
        by = int(float(getattr(el, "y", 0.0)) * height)
        bw = max(1, int(float(getattr(el, "w", 0.2)) * width))
        bh = max(1, int(float(getattr(el, "h", 0.2)) * height))
        fitted = ImageOps.contain(img, (bw, bh))
        ox = bx + (bw - fitted.width) // 2
        oy = by + (bh - fitted.height) // 2
        base_rgba.alpha_composite(fitted, (ox, oy))
    except Exception:
        pass
    finally:
        if tmp_frame is not None:
            try:
                tmp_frame.unlink(missing_ok=True)
            except Exception:
                pass


def render_card_to_png(
    card: TitleCard,
    out_path: Path,
    width: int = 1920,
    height: int = 1080,
) -> bool:
    """Render a card — solid-color / image / video background plus its
    styled title + subtitle (with box / outline / shadow effects) — to
    a STILL PNG.

    A "designed" content slide (PowerPoint-style: pick a background,
    type text, style it) stores this PNG as its ``image_path`` so it
    flows through EVERY existing pipeline — deck stitch, xfade
    transitions, PowerPoint export, list thumbnails — with no
    card-specific handling anywhere. The page keeps its ``card`` too so
    the writer can re-open the designer and re-render. Returns True on
    success."""
    try:
        from PIL import Image, ImageOps
    except Exception:
        return False
    out_path.parent.mkdir(parents=True, exist_ok=True)
    kind = (getattr(card, "kind", "color") or "color").lower()
    media = getattr(card, "bg_media_path", "") or ""
    base = None
    if kind == "image" and media and Path(media).exists():
        try:
            src = Image.open(media).convert("RGB")
            base = ImageOps.fit(
                src, (width, height), Image.LANCZOS)
        except Exception:
            base = None
    elif kind == "video" and media and Path(media).exists():
        # Best-effort: grab the first frame so a still exists.
        try:
            frame = out_path.with_name(out_path.stem + "_frame.png")
            subprocess.run(
                ["ffmpeg", "-y", "-i",
                 str(Path(media).resolve()),
                 "-vframes", "1", str(frame)],
                capture_output=True, timeout=30)
            if frame.exists():
                src = Image.open(str(frame)).convert("RGB")
                base = ImageOps.fit(
                    src, (width, height), Image.LANCZOS)
                frame.unlink(missing_ok=True)
        except Exception:
            base = None
    if base is None:
        rgba = _hex_to_rgba(
            getattr(card, "bg_color", "#000000"), "#000000")
        base = Image.new("RGB", (width, height), rgba[:3])
    # ── Composite free-placed elements (image / video-still / text)
    #    in LAYER order so what the designer shows matches the still. ──
    from PIL import ImageDraw
    base_rgba = base.convert("RGBA")
    draw = ImageDraw.Draw(base_rgba)
    scale = height / 1080.0
    fontfile = _resolve_font_file()
    for el in sorted(
            (getattr(card, "elements", None) or []),
            key=lambda e: getattr(e, "z", 0)):
        kind = getattr(el, "kind", "text")
        if kind in ("image", "video"):
            _paste_media_element(base_rgba, el, width, height, out_path)
        elif kind == "text" and (
                getattr(el, "text", "") or "").strip() and fontfile:
            _draw_text_box(draw, el, width, height, scale, fontfile)
    # Legacy title / subtitle block on top (title / ending cards; a
    # designed slide leaves these blank). Render via a copy with the
    # elements stripped so it draws ONLY the title/subtitle.
    if (getattr(card, "title", "") or "").strip() or (
            getattr(card, "subtitle", "") or "").strip():
        try:
            tmp = card.model_copy(deep=True)
            tmp.elements = []
            tmp.text_boxes = []
            txt_png = out_path.with_name(out_path.stem + "_txt.png")
            if _render_text_overlay_png(tmp, width, height, txt_png):
                txt = Image.open(str(txt_png)).convert("RGBA")
                base_rgba.alpha_composite(txt)
                txt_png.unlink(missing_ok=True)
        except Exception:
            pass
    base = base_rgba.convert("RGB")
    try:
        base.save(str(out_path))
        return True
    except Exception:
        return False


def render_card_to_mp4(
    card: TitleCard,
    output_path: Path,
    duration: float,
    audio_path: str = "",
    width: int = 1280,
    height: int = 720,
    fps: int = 30,
) -> Tuple[bool, str]:
    """Render a title / ending card to an MP4 of ``duration``
    seconds: a color / image / video background with the card's
    styled title + subtitle over it (fading in and out), plus the
    card's audio when supplied. The output concatenates with the
    other deck segments."""
    if not ffmpeg_available():
        return (False, "ffmpeg not found on PATH.")
    dur = max(0.5, float(duration or 0.5))
    output_path.parent.mkdir(parents=True, exist_ok=True)
    kind = (getattr(card, "kind", "color") or "color").lower()
    media = getattr(card, "bg_media_path", "") or ""
    inputs: List[str] = []
    filters: List[str] = []
    scale_crop = (
        f"scale={width}:{height}:"
        f"force_original_aspect_ratio=increase,"
        f"crop={width}:{height},setsar=1,fps={fps}")
    if kind == "image" and media and Path(media).exists():
        inputs += [
            "-loop", "1", "-t", f"{dur:.3f}",
            "-i", str(Path(media).resolve())]
        filters.append(f"[0:v]{scale_crop},format=yuv420p[bg]")
    elif kind == "video" and media and Path(media).exists():
        inputs += ["-stream_loop", "-1",
                   "-i", str(Path(media).resolve())]
        filters.append(
            f"[0:v]{scale_crop},trim=0:{dur:.3f},"
            f"setpts=PTS-STARTPTS,format=yuv420p[bg]")
    else:
        color = _sanitize_color(getattr(card, "bg_color", "#000000"))
        inputs += [
            "-f", "lavfi", "-t", f"{dur:.3f}",
            "-i", f"color=c={color}:s={width}x{height}:r={fps}"]
        filters.append("[0:v]format=yuv420p[bg]")
    n_video_inputs = 1  # the background is input 0
    last = "bg"
    # --- free-placed image / video ELEMENTS, layered by z under the
    #     text (each scaled to fit its box, centered via pad; videos
    #     loop to fill the slide and play muted) ---
    media_els = sorted(
        [e for e in (getattr(card, "elements", None) or [])
         if getattr(e, "kind", "") in ("image", "video")
         and getattr(e, "media_path", "")
         and Path(e.media_path).exists()],
        key=lambda e: getattr(e, "z", 0))
    for el in media_els:
        ex = int(float(el.x) * width)
        ey = int(float(el.y) * height)
        ew = max(2, int(float(el.w) * width))
        eh = max(2, int(float(el.h) * height))
        idx = n_video_inputs
        fit = (
            f"scale={ew}:{eh}:force_original_aspect_ratio=decrease,"
            f"pad={ew}:{eh}:(ow-iw)/2:(oh-ih)/2:color=black@0")
        if el.kind == "image":
            inputs += [
                "-loop", "1", "-t", f"{dur:.3f}",
                "-i", str(Path(el.media_path).resolve())]
            filters.append(
                f"[{idx}:v]{fit},format=rgba[e{idx}]")
        else:
            inputs += [
                "-stream_loop", "-1",
                "-i", str(Path(el.media_path).resolve())]
            filters.append(
                f"[{idx}:v]{fit},trim=0:{dur:.3f},"
                f"setpts=PTS-STARTPTS,format=rgba[e{idx}]")
        filters.append(
            f"[{last}][e{idx}]overlay={ex}:{ey}[m{idx}]")
        last = f"m{idx}"
        n_video_inputs += 1
    # --- styled text as a faded overlay (PIL PNG, no drawtext) ---
    fade = max(0.0, float(getattr(card, "text_fade_seconds", 0.0)
                          or 0.0))
    text_png = output_path.with_name(
        output_path.stem + "_cardtext.png")
    made_text = _render_text_overlay_png(
        card, width, height, text_png)
    if made_text:
        inputs += [
            "-loop", "1", "-t", f"{dur:.3f}",
            "-i", str(text_png.resolve())]
        txt_idx = n_video_inputs
        n_video_inputs += 1
        fchain = "format=rgba"
        if fade > 0:
            out_st = max(0.0, dur - fade)
            fchain += (
                f",fade=t=in:st=0:d={fade:.3f}:alpha=1"
                f",fade=t=out:st={out_st:.3f}:d={fade:.3f}:alpha=1")
        filters.append(f"[{txt_idx}:v]{fchain}[txt]")
        filters.append(f"[{last}][txt]overlay=0:0[v]")
        last = "v"
    filter_str = ";".join(filters)
    has_audio = bool(
        audio_path and Path(audio_path).exists()
        and Path(audio_path).stat().st_size > 0)
    if has_audio:
        inputs += ["-i", str(Path(audio_path).resolve())]
    cmd = [
        "ffmpeg", "-y",
        *inputs,
        "-filter_complex", filter_str,
        "-map", f"[{last}]",
    ]
    if has_audio:
        cmd += ["-map", f"{n_video_inputs}:a",
                "-c:a", "aac", "-profile:a", "aac_low",
                "-ar", "48000", "-b:a", "192k", "-ac", "2"]
    cmd += [
        "-c:v", "libx264", "-pix_fmt", "yuv420p", "-r", str(fps),
        "-t", f"{dur:.3f}",
        "-movflags", "+faststart",
        str(output_path.resolve()),
    ]
    try:
        proc = subprocess.run(
            cmd, capture_output=True, text=True, timeout=600)
    except subprocess.TimeoutExpired:
        return (False, "card render timed out.")
    finally:
        try:
            text_png.unlink(missing_ok=True)
        except Exception:
            pass
    if proc.returncode != 0:
        return (False,
                "ffmpeg card render failed. stderr (last 400):\n"
                + (proc.stderr or "")[-400:])
    return (True, f"Card rendered to {output_path}.")


def render_black_spacer(
    output_path: Path,
    duration: float,
    width: int = 1280,
    height: int = 720,
    fps: int = 30,
) -> Tuple[bool, str]:
    """Render a silent BLACK video of ``duration`` seconds — the
    dark-space scene break inserted between groups. No audio track
    (the concat synthesizes silence / the bed carries over)."""
    if not ffmpeg_available():
        return (False, "ffmpeg not found on PATH.")
    dur = max(0.1, float(duration or 0.0))
    output_path.parent.mkdir(parents=True, exist_ok=True)
    cmd = [
        "ffmpeg", "-y",
        "-f", "lavfi", "-t", f"{dur:.3f}",
        "-i", f"color=c=black:s={width}x{height}:r={fps}",
        "-c:v", "libx264", "-pix_fmt", "yuv420p", "-r", str(fps),
        "-t", f"{dur:.3f}", "-movflags", "+faststart",
        str(output_path.resolve()),
    ]
    try:
        proc = subprocess.run(
            cmd, capture_output=True, text=True, timeout=120)
    except subprocess.TimeoutExpired:
        return (False, "black spacer render timed out.")
    if proc.returncode != 0:
        return (False,
                "black spacer render failed:\n"
                + (proc.stderr or "")[-300:])
    return (True, f"Black spacer ({dur:.2f}s) rendered.")


def ordered_pages(deck: SlideDeckProject) -> List[SlidePage]:
    """Pages in the deck's ACTUAL group / render order: walk
    ``deck.groups`` in order, and within each group its member slides
    sorted by timeline placement (placed first, by start time; then
    tray slides by index); finally any orphan pages (no group, or a
    group not in ``deck.groups``) in their existing order.

    The slide list, the agent's slide numbering, and the render all
    use this so "slide #N" means the same thing everywhere — the left
    list matches the group order shown in the center."""
    out: List[SlidePage] = []
    seen: set = set()
    group_ids = {g.id for g in (deck.groups or [])}
    for g in (deck.groups or []):
        members = [
            p for p in deck.pages
            if getattr(p, "group_id", None) == g.id]
        members.sort(key=lambda p: (
            getattr(p, "start_time_seconds_in_group", None) is None,
            float(getattr(p, "start_time_seconds_in_group", 0.0)
                  or 0.0),
            int(getattr(p, "index", 0) or 0)))
        for p in members:
            out.append(p)
            seen.add(p.id)
    for p in deck.pages:
        gid = getattr(p, "group_id", None)
        if p.id not in seen and (not gid or gid not in group_ids):
            out.append(p)
            seen.add(p.id)
    # Safety: include anything somehow missed (e.g. a page whose
    # group_id points at a group not walked above) so nothing vanishes.
    for p in deck.pages:
        if p.id not in seen:
            out.append(p)
    return out


def group_card_page(
    deck: SlideDeckProject, group: SlideGroup,
) -> Optional[SlidePage]:
    """When ``group`` is a CARD group — its single placed slide is a
    title / ending card — return that card page; otherwise None. A
    card lives in its own group so it gets the same audio overlay
    and inter-group transition as any other group."""
    placed = [
        p for p in deck.pages
        if getattr(p, "group_id", None) == group.id
        and getattr(p, "start_time_seconds_in_group", None)
        is not None]
    if len(placed) == 1 and getattr(placed[0], "card", None):
        return placed[0]
    return None


def render_card_group_to_mp4(
    deck: SlideDeckProject,
    group: SlideGroup,
    card_page: SlidePage,
    output_path: Path,
    width: int = 1280,
    height: int = 720,
    fps: int = 30,
) -> Tuple[bool, str]:
    """Render a card group: the card visual (via
    ``render_card_to_mp4``) held for the card's duration — stretched
    to the group's audio overlay when that runs longer — with the
    group's composed narration under it."""
    overlay = resolve_group_overlay(group, working_dir=deck.working_dir)
    overlay_dur = float(
        getattr(group, "overlay_audio_duration_seconds", 0.0) or 0.0)
    dur = max(0.5, float(
        getattr(card_page, "duration_seconds", 0.0) or 0.0))
    if overlay and Path(overlay).exists() and overlay_dur > dur:
        dur = overlay_dur
    return render_card_to_mp4(
        card_page.card, output_path, dur,
        audio_path=(overlay or ""),
        width=width, height=height, fps=fps)


def render_group_to_mp4(
    deck: SlideDeckProject,
    group: SlideGroup,
    output_path: Path,
    width: int = 1280,
    height: int = 720,
    fps: int = 30,
    exclude_background: bool = False,
) -> Tuple[bool, str]:
    """Render ONE group into an MP4 the same way the group
    editor's "Preview" button does.

    Pipeline:
      * Pull the group's placed slides (those with
        ``start_time_seconds_in_group`` set), sort by start
        time.
      * Each slide's render hold = gap to next placed slide;
        the LAST slide stretches to the overlay end if that's
        longer than its own duration.
      * Strip per-slide audio_path (the overlay owns audio
        for the group's span).
      * Attach the group's composed overlay to the first
        slide so the stitcher's ``adelay`` lands it at deck-
        time 0 for this group.
      * Hand the synthetic deck to
        ``stitch_slide_deck_to_mp4``.

    This helper is the single source of truth for "what a
    group sounds + looks like" — both the group editor's
    preview and the slide deck editor's concat-based deck
    preview call it so they stay byte-for-byte consistent.
    """
    placed = sorted(
        (p for p in deck.pages
         if p.group_id == group.id
         and getattr(
             p, "start_time_seconds_in_group", None)
         is not None),
        key=lambda p: float(
            getattr(
                p, "start_time_seconds_in_group", 0.0)
            or 0.0))
    if not placed:
        # A group can carry a fully recorded + edited audio track
        # while its member slides still sit in the tray (the writer
        # recorded narration but never dragged the slides onto the
        # group timeline). Skipping the group outright then drops
        # its sound from every export — the deck renders its
        # visuals as orphan slides but comes out silent, which is
        # exactly the "12 minutes of video, no audio" symptom.
        #
        # Recover gracefully: if the group has audio and at least
        # one member slide, auto-place the members evenly across
        # the narration so the sound has slides to ride on. Uses
        # deep-copied pages so the writer's saved placement (empty)
        # is not mutated behind their back.
        overlay_probe = resolve_group_overlay(
            group, working_dir=deck.working_dir)
        if overlay_probe and Path(overlay_probe).exists():
            members = [
                p for p in deck.pages
                if p.group_id == group.id]
            order = {
                pid: i for i, pid in enumerate(
                    getattr(group, "page_ids", []) or [])}
            members.sort(
                key=lambda p: (
                    order.get(p.id, 1_000_000),
                    getattr(p, "index", 0)))
            if members:
                span = float(
                    getattr(
                        group,
                        "overlay_audio_duration_seconds",
                        0.0) or 0.0)
                if span <= 0:
                    span = max(1.0, len(members) * 3.0)
                per = max(0.25, span / len(members))
                placed = []
                for i, m in enumerate(members):
                    c = m.model_copy(deep=False)
                    c.start_time_seconds_in_group = round(
                        i * per, 3)
                    placed.append(c)
                print(
                    f"[slide_deck] group "
                    f"'{group.name or group.id}' had audio but "
                    f"no placed slides — auto-placed "
                    f"{len(placed)} member slide(s) across "
                    f"{span:.1f}s so its narration exports.")
        if not placed:
            return (False,
                    f"Group '{group.name or group.id}' has no "
                    "placed slides.")
    overlay_dur = float(
        getattr(
            group,
            "overlay_audio_duration_seconds", 0.0) or 0.0)
    render_pages: list = []
    for i, src in enumerate(placed):
        cur_start = float(
            getattr(
                src, "start_time_seconds_in_group", 0.0)
            or 0.0)
        if i + 1 < len(placed):
            next_start = float(
                getattr(
                    placed[i + 1],
                    "start_time_seconds_in_group", 0.0)
                or 0.0)
            hold = max(0.25, next_start - cur_start)
        else:
            own = max(
                0.25, float(
                    getattr(
                        src, "duration_seconds", 0.0)
                    or 0.0))
            tail = max(0.0, overlay_dur - cur_start)
            if tail > 0:
                # There is narration under this group. The last
                # slide should stretch to cover it, but NOT sit
                # frozen and silent for long after the voice stops
                # — that trailing dead-air reads as a ~0.5s gap
                # before the next group on a cut (a transition
                # hides it by overlapping, which is why the gap
                # only shows without one). Cap the post-narration
                # hold to a few frames so groups flow straight
                # into each other while a transition, if set, still
                # has content to fade over.
                hold = max(
                    tail,
                    min(own, tail + MAX_TRAILING_HOLD_SECONDS))
            else:
                hold = own
        copy = src.model_copy(deep=False)
        copy.duration_seconds = round(hold, 3)
        render_pages.append(copy)
    # Resolve (and, if the cache is missing, recompose) the
    # group's overlay so every audio-clip edit — gain, de-essing,
    # noise reduction, looping — is baked into the render even
    # when the cached WAV went missing. When ``exclude_background``
    # is set (a universal deck bed is replacing group beds), render
    # the group's narration WITHOUT its own background-loop lanes so
    # the two beds don't fight.
    overlay_path = compose_group_overlay_variant(
        group, deck.working_dir, exclude_background)
    has_overlay = bool(
        overlay_path and Path(overlay_path).exists())
    if has_overlay:
        # The composed overlay is the group's authoritative audio
        # — it already blends every take + edit — so per-slide
        # audio_path values are redundant and would double up.
        # Strip them and let the overlay own the group's sound.
        for copy in render_pages:
            copy.audio_path = ""
        # Attach the overlay to the first page that has a VALID
        # image on disk. ``stitch_slide_deck_to_mp4`` drops any
        # page whose image is missing before it builds the audio
        # mix — so if we blindly hung the overlay on
        # ``render_pages[0]`` and that first slide's image happened
        # to be gone, the audio was silently discarded along with
        # it while the remaining slides still rendered. That is the
        # "images export, audio doesn't" symptom. Anchoring the
        # audio to a surviving page guarantees the mix step sees
        # it.
        anchor = next(
            (c for c in render_pages
             if getattr(c, "image_path", "")
             and Path(c.image_path).exists()),
            render_pages[0])
        anchor.audio_path = overlay_path
    # else: NO group overlay — keep each slide's own audio_path so
    # narration the writer attached per-slide (via the main slide
    # editor's Import / Record / Edit) still exports. Stripping it
    # here is what silently muted decks whose audio lived on the
    # slides rather than in a group audio track.
    synthetic = SlideDeckProject(
        id=f"render_{group.id}",
        name=f"Group {group.name or group.id}",
        working_dir=deck.working_dir,
        pages=render_pages,
    )
    return stitch_slide_deck_to_mp4(
        synthetic, output_path,
        width=width, height=height, fps=fps)


def stitch_slide_deck_to_mp4(
    deck: SlideDeckProject,
    output_path: Path,
    width: int = 1280,
    height: int = 720,
    fps: int = 30,
) -> Tuple[bool, str]:
    """Render the deck into a single MP4: each page is an image
    held for its duration with its audio mixed in.

    Returns ``(success, message)``. Uses the existing
    ``stitch_clips`` for the visual concat, then ffmpeg with
    filter_complex to attach per-slide audio offsets.
    """
    if not ffmpeg_available():
        return (
            False,
            "ffmpeg not found on PATH. Install ffmpeg and try again.")
    pages = [
        p for p in deck.pages
        if p.image_path and Path(p.image_path).exists()]
    if not pages:
        return (False, "No slides with usable images.")
    output_path.parent.mkdir(parents=True, exist_ok=True)
    # Bake any per-slide TEXT OVERLAY onto its image (into a temp
    # file) so the rest of the pipeline treats it as a normal still.
    # Non-mutating — we only swap the path used for this render.
    _baked: List[Path] = []
    image_paths: List[Path] = []
    for p in pages:
        src_img = Path(p.image_path)
        ov = getattr(p, "text_overlay", None)
        if ov is not None:
            try:
                baked = output_path.parent / (
                    f"_ov_{p.id}"
                    f"{Path(p.image_path).suffix or '.png'}")
                if bake_text_overlay(p.image_path, ov, baked):
                    src_img = baked
                    _baked.append(baked)
            except Exception:
                pass
        image_paths.append(src_img)
    image_durations = [
        max(MIN_SLIDE_SECONDS, float(p.duration_seconds))
        for p in pages]
    output_path.parent.mkdir(parents=True, exist_ok=True)
    silent_path = output_path.with_name(
        output_path.stem + "_silent.mp4")
    # When any slide has a non-cut transition, route through the
    # transition stitcher (xfade) so writers' picks land in the
    # final render. Otherwise stick with the simpler concat path.
    has_transitions = any(
        (p.transition_in or "cut") != "cut"
        and float(p.transition_seconds or 0.0) > 0
        for p in pages[1:])
    if has_transitions:
        transitions = [
            (p.transition_in or "cut",
             float(p.transition_seconds or 0.0))
            for p in pages]
        visual_result = stitch_with_transitions(
            image_paths, silent_path,
            clip_durations=image_durations,
            transitions=transitions,
            width=width, height=height, fps=fps)
        # Each non-cut transition compresses the timeline by
        # ``transition_seconds`` — adjust the per-slide start
        # offsets we use later for audio mixing.
        offsets_adjusted = True
    else:
        visual_result = stitch_clips(
            image_paths, silent_path,
            clip_durations=image_durations)
        offsets_adjusted = False
    # The baked overlay stills have been consumed by the visual
    # stitch — drop them now.
    for bp in _baked:
        try:
            bp.unlink(missing_ok=True)
        except Exception:
            pass
    if not visual_result.success:
        return (False, visual_result.error)
    # No audio at all? Done.
    audio_pages = [
        (i, p) for i, p in enumerate(pages)
        if p.audio_path and Path(p.audio_path).exists()
        and Path(p.audio_path).stat().st_size > 0
    ]
    if not audio_pages:
        try:
            silent_path.rename(output_path)
        except Exception as e:
            return (
                False,
                f"Could not finalize file: {e}")
        return (True, f"Slide deck saved to {output_path}.")
    # Compute per-page start offsets so each audio take lines up
    # with its slide. When transitions are present, each non-cut
    # boundary's seconds shorten the cumulative timeline (xfade
    # overlaps the prior clip's tail with the next clip's head).
    starts: List[float] = []
    running = 0.0
    for i, d in enumerate(image_durations):
        starts.append(running)
        running += d
        if has_transitions and i + 1 < len(pages):
            next_page = pages[i + 1]
            kind = next_page.transition_in or "cut"
            secs = float(next_page.transition_seconds or 0.0)
            if kind != "cut" and secs > 0:
                running -= secs
    # Build the audio mix filter.
    inputs: List[str] = ["-i", str(silent_path.resolve())]
    filter_parts: List[str] = []
    mix_labels: List[str] = []
    for idx, (page_idx, page) in enumerate(audio_pages, start=1):
        inputs.extend(
            ["-i", str(Path(page.audio_path).resolve())])
        start = starts[page_idx]
        chain = ["asetpts=PTS-STARTPTS"]
        if start > 0:
            chain.append(
                f"adelay={int(start * 1000)}|{int(start * 1000)}")
        chain.append("aformat=channel_layouts=stereo")
        label = f"a{idx}"
        filter_parts.append(
            f"[{idx}:a]" + ",".join(chain) + f"[{label}]")
        mix_labels.append(f"[{label}]")
    filter_parts.append(
        "".join(mix_labels)
        + f"amix=inputs={len(mix_labels)}:duration=longest:"
          "normalize=0[mix]")
    filter_str = ";".join(filter_parts)
    # The visual track is authoritative — we want the full slide
    # run-time even when a slide has shorter (or no) audio. Using
    # ``-shortest`` here would cut the deck the moment the last
    # audio take ends, dropping any trailing silent slides.
    cmd = [
        "ffmpeg", "-y",
        *inputs,
        "-filter_complex", filter_str,
        "-map", "0:v:0", "-map", "[mix]",
        "-c:v", "copy",
        "-c:a", "aac",
        "-t", f"{running:.3f}",
        str(output_path),
    ]
    try:
        proc = subprocess.run(
            cmd, capture_output=True, text=True, timeout=900)
        if proc.returncode != 0:
            return (
                False,
                "ffmpeg slide mux failed. stderr (last 400):\n"
                + (proc.stderr or "")[-400:])
    except subprocess.TimeoutExpired:
        return (False, "ffmpeg slide mux timed out (15 min).")
    finally:
        try:
            silent_path.unlink(missing_ok=True)
        except Exception:
            pass
    return (True, f"Slide deck saved to {output_path}.")
