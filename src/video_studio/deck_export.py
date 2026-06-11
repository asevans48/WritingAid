"""Chapter-level deck export.

Walks a chapter's scenes in order, picks each scene's chosen output
(favorite image, video clip, or slideshow stitch), optionally renders
a title card per scene, and concatenates the result into a single
MP4 the writer can share — a "slide deck of the chapter."

The heavy lifting (mixed image+video concatenation) lives in
``stitcher.stitch_clips``; this module just builds the list of
(path, duration) tuples and renders the title cards.
"""

from __future__ import annotations

from pathlib import Path
from typing import Any, List, Optional, Tuple

from PIL import Image, ImageDraw, ImageFont


# Card geometry matches the stitcher's default 30 fps / yuv420p
# pipeline; 1280x720 looks clean alongside 16:9 video clips and
# scales nicely for both phones and laptops.
TITLE_CARD_SIZE = (1280, 720)
TITLE_CARD_DURATION_SECONDS = 2.5


def render_title_card(
    out_path: Path,
    title: str,
    subtitle: str = "",
    overline: str = "",
    size: Tuple[int, int] = TITLE_CARD_SIZE,
) -> Path:
    """Render a clean, dark-background title card with one scene's
    name + optional metadata. PNG, ready for the stitcher.

    Layout::

        ┌────────────────────────────────────────┐
        │ <overline (small caps, dim)>           │
        │                                        │
        │         <title (large, bold)>          │
        │                                        │
        │ <subtitle (medium, dim)>               │
        └────────────────────────────────────────┘
    """
    w, h = size
    img = Image.new("RGB", size, color=(15, 23, 42))  # slate-900
    draw = ImageDraw.Draw(img)

    # Fonts — fall back to PIL's default when the system fonts
    # we'd prefer aren't reachable. Default is small and ugly but
    # always present, which beats crashing the export.
    title_font = _load_font(72, bold=True)
    subtitle_font = _load_font(28)
    overline_font = _load_font(20)

    pad = 80

    # Overline (small caps look, on a dim color).
    if overline:
        over_text = overline.upper()
        draw.text(
            (pad, pad),
            over_text,
            fill=(148, 163, 184),  # slate-400
            font=overline_font)

    # Title — wrap manually so very long names don't overflow.
    title_lines = _wrap_text(title, title_font, w - 2 * pad, draw)
    # Vertically center the title block.
    line_h = title_font.size + 12 if hasattr(title_font, "size") else 80
    total_h = line_h * len(title_lines)
    title_top = (h - total_h) // 2 - 20
    for i, line in enumerate(title_lines):
        draw.text(
            (pad, title_top + i * line_h),
            line,
            fill=(248, 250, 252),  # slate-50
            font=title_font)

    # Subtitle — under the title block.
    if subtitle:
        sub_y = title_top + total_h + 40
        sub_lines = _wrap_text(
            subtitle, subtitle_font, w - 2 * pad, draw)
        sub_line_h = (
            subtitle_font.size + 8
            if hasattr(subtitle_font, "size") else 32)
        for i, line in enumerate(sub_lines):
            draw.text(
                (pad, sub_y + i * sub_line_h),
                line,
                fill=(148, 163, 184),
                font=subtitle_font)

    img.save(out_path, "PNG")
    return out_path


def _load_font(size: int, bold: bool = False) -> Any:
    """Try a list of common system fonts; fall back to PIL default
    when nothing matches. PIL's default font ignores size, so the
    output looks bad — but it never crashes, which is what we want
    for export."""
    candidates = (
        # macOS — modern + classic locations.
        ("/System/Library/Fonts/SFNS.ttf"
         if not bold else
         "/System/Library/Fonts/SFNSMono.ttf"),
        "/System/Library/Fonts/Helvetica.ttc",
        "/System/Library/Fonts/Supplemental/Arial.ttf",
        # Linux.
        ("/usr/share/fonts/truetype/dejavu/DejaVuSans-Bold.ttf"
         if bold else
         "/usr/share/fonts/truetype/dejavu/DejaVuSans.ttf"),
        # Windows.
        ("C:/Windows/Fonts/arialbd.ttf"
         if bold else
         "C:/Windows/Fonts/arial.ttf"),
    )
    for path in candidates:
        try:
            return ImageFont.truetype(path, size)
        except (OSError, IOError):
            continue
    return ImageFont.load_default()


def _wrap_text(text: str, font: Any, max_w: int, draw) -> List[str]:
    """Greedy word-wrap into lines that fit ``max_w`` pixels. No
    hyphenation — long URLs / single tokens may overflow but won't
    crash the layout."""
    if not text:
        return []
    words = text.split()
    if not words:
        return []
    lines: List[str] = []
    current: List[str] = []
    for word in words:
        trial = " ".join(current + [word])
        try:
            bbox = draw.textbbox((0, 0), trial, font=font)
            width = bbox[2] - bbox[0]
        except Exception:
            # textbbox missing on very old PIL; fall back to length.
            width = len(trial) * 10
        if width <= max_w or not current:
            current.append(word)
        else:
            lines.append(" ".join(current))
            current = [word]
    if current:
        lines.append(" ".join(current))
    return lines


def collect_chapter_scenes(
    studio: Any, chapter_id: str,
) -> List[Any]:
    """Return the scenes that belong to a given chapter, ordered for
    deck export.

    Ordering precedence:
      1. Topological order from the chapter's first scene if the
         chapter's scenes are connected by hops.
      2. Otherwise grid reading order (top-to-bottom, left-to-right).
    """
    if studio is None or not chapter_id:
        return []
    members = [
        s for s in studio.scenes
        if getattr(s, "chapter_id", None) == chapter_id
    ]
    if not members:
        return []
    # Try hops first — if any scene in the chapter has incoming or
    # outgoing hops, prefer that order. Use studio's BFS helper.
    try:
        if any(
            (h.from_scene_id in {s.id for s in members}
             or h.to_scene_id in {s.id for s in members})
            for h in (getattr(studio, "hops", []) or [])
        ):
            # Pick the start: first chapter scene with NO incoming
            # hop within this chapter (so we start at the head).
            chapter_ids = {s.id for s in members}
            incoming = {
                h.to_scene_id for h in studio.hops
                if h.from_scene_id in chapter_ids
            }
            heads = [s for s in members if s.id not in incoming]
            start = heads[0].id if heads else members[0].id
            ordered_ids = [
                s.id for s in
                studio.topological_order_starting_at(start)
                if s.id in chapter_ids
            ]
            # If BFS missed any (disconnected components), append in
            # grid order so nothing is silently dropped.
            seen = set(ordered_ids)
            fallback = sorted(
                (s for s in members if s.id not in seen),
                key=lambda s: (s.grid_row, s.grid_col))
            id_to_scene = {s.id: s for s in members}
            return [id_to_scene[i] for i in ordered_ids] + fallback
    except Exception:
        pass
    # Grid reading order — top to bottom, left to right.
    return sorted(
        members, key=lambda s: (s.grid_row, s.grid_col))


def export_chapter_pptx(
    scenes: List[Any],
    output_path: Path,
    chapter_title: str = "",
) -> Tuple[bool, str, List[str]]:
    """Compose a PowerPoint (.pptx) from the chapter's action images.

    One slide per ACTION favorite image (or per scene when the
    scene has no actions). Slides are intentionally empty — no
    titles, no descriptions, no overlays — so the writer can take
    the deck into PowerPoint / Keynote / Slides and arrange,
    annotate, or re-time freely without first deleting our text.
    Each image is fitted to the slide preserving its aspect ratio
    and centered on a black background.

    Selection rule per scene:
      * Slideshow scenes → walk every action and use
        ``action.favorite_image()`` (the favorite when starred,
        or the first generated image as fallback).
      * Non-slideshow scenes (single image or video mode) → use
        ``scene.favorite_clip()`` so the writer still gets one
        slide per scene to insert / replace later.

    Returns ``(success, message, skipped)``. ``skipped`` is a list
    of human-readable per-action / per-scene reasons.

    python-pptx is an optional dependency. When unavailable, the
    return value is ``(False, "python-pptx not installed: …", [])``
    so the host can surface a clear install hint.
    """
    try:
        from pptx import Presentation
        from pptx.util import Emu, Inches
        from pptx.dml.color import RGBColor
    except Exception as e:
        return (
            False,
            (
                "python-pptx isn't installed in this environment. "
                "Install it with:\n  pip install python-pptx\n\n"
                f"Underlying error: {e}"),
            [])
    if not scenes:
        return (
            False, "No scenes to export.", [])
    # 16:9 widescreen — the universal slide deck aspect ratio.
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)
    blank_layout = prs.slide_layouts[6]
    slide_w_emu = int(prs.slide_width)
    slide_h_emu = int(prs.slide_height)
    skipped: List[str] = []

    IMAGE_EXTS = {
        ".png", ".jpg", ".jpeg", ".webp", ".gif", ".bmp"}
    VIDEO_EXTS = {
        ".mp4", ".mov", ".webm", ".mkv", ".avi", ".m4v"}

    def _add_image_slide(
        image_path: Path, label: str,
    ) -> None:
        """Drop one full-bleed image slide on a black background.
        No text — the deck is a sequence of images the writer can
        edit / re-arrange directly in PowerPoint."""
        slide = prs.slides.add_slide(blank_layout)
        _set_slide_background(slide, RGBColor(0x00, 0x00, 0x00))
        try:
            left, top, width, height = _fit_to_slide(
                image_path, slide_w_emu, slide_h_emu)
            slide.shapes.add_picture(
                str(image_path), left, top,
                width=width, height=height)
        except Exception as e:
            skipped.append(
                f"{label}: PPTX embed failed ({e})")

    def _add_video_slide(
        video_path: Path, label: str,
    ) -> None:
        """Drop one full-bleed embedded video on a black slide.
        The poster frame (when ffmpeg is available) shows up
        before play; otherwise PowerPoint renders an empty
        placeholder until the writer clicks."""
        slide = prs.slides.add_slide(blank_layout)
        _set_slide_background(slide, RGBColor(0x00, 0x00, 0x00))
        try:
            poster = _make_video_poster(video_path)
            if poster is not None:
                left, top, width, height = _fit_to_slide(
                    poster, slide_w_emu, slide_h_emu)
            else:
                left = top = Emu(0)
                width = slide_w_emu
                height = slide_h_emu
            slide.shapes.add_movie(
                str(video_path), left, top, width, height,
                poster_frame_image=(
                    str(poster) if poster else None))
        except Exception as e:
            skipped.append(
                f"{label}: PPTX embed failed ({e})")

    def _emit_clip_slide(
        clip_path: Path, label: str,
    ) -> None:
        suffix = clip_path.suffix.lower()
        if suffix in IMAGE_EXTS:
            _add_image_slide(clip_path, label)
        elif suffix in VIDEO_EXTS:
            _add_video_slide(clip_path, label)
        else:
            skipped.append(
                f"{label}: unsupported format ({suffix})")

    def _media_is_usable(path_str: str) -> Tuple[bool, str]:
        if not path_str:
            return False, "empty path"
        try:
            p = Path(path_str)
        except Exception:
            return False, "invalid path"
        if not p.exists():
            return False, f"file not found ({p.name})"
        try:
            if p.stat().st_size == 0:
                return False, f"file is 0 bytes ({p.name})"
        except Exception as e:
            return False, f"stat error ({e})"
        return True, ""

    for idx, scene in enumerate(scenes, start=1):
        scene_label = scene.name or f"Scene {idx}"
        actions = (
            getattr(scene, "actions", None) or []
            if (getattr(scene, "mode", "") == "slideshow")
            else [])
        if actions:
            # One slide per action — favorite image first, then
            # fall back to images[0] via favorite_image().
            for a_idx, action in enumerate(actions, start=1):
                a_label = (
                    f"{scene_label} → "
                    + (action.name or f"action {a_idx}"))
                chosen = action.favorite_image()
                if chosen is None:
                    skipped.append(
                        f"{a_label}: no images on action")
                    continue
                ok, why = _media_is_usable(chosen.file_path)
                if not ok:
                    skipped.append(f"{a_label}: {why}")
                    continue
                if getattr(chosen, "is_placeholder", False):
                    skipped.append(
                        f"{a_label}: placeholder only")
                _add_image_slide(
                    Path(chosen.file_path), a_label)
            continue
        # Non-slideshow scene — use the scene's favorite clip so
        # the writer still gets one slide per scene for re-arranging.
        clip = scene.favorite_clip()
        if clip is None:
            skipped.append(f"{scene_label}: no favorite output")
            continue
        ok, why = _media_is_usable(clip.file_path)
        if not ok:
            skipped.append(f"{scene_label}: {why}")
            continue
        if getattr(clip, "is_placeholder", False):
            skipped.append(f"{scene_label}: placeholder only")
        _emit_clip_slide(Path(clip.file_path), scene_label)

    if len(prs.slides) == 0:
        return (
            False,
            "No usable images / clips in this chapter. "
            "Generate or upload images for the scene actions "
            "first, then try again.",
            skipped)
    output_path.parent.mkdir(parents=True, exist_ok=True)
    try:
        prs.save(str(output_path))
    except Exception as e:
        return (False, f"PowerPoint save failed: {e}", skipped)
    return (
        True,
        f"PowerPoint deck saved to {output_path}.",
        skipped)


def _set_slide_background(slide, rgb_color) -> None:
    """Set a solid-color background on a single slide. python-pptx
    doesn't expose this directly so we drill through the XML."""
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = rgb_color


def _fit_to_slide(
    image_path: Path, slide_w_emu: int, slide_h_emu: int,
) -> Tuple[int, int, int, int]:
    """Compute (left, top, width, height) in EMUs that fits the
    image inside the slide while preserving its aspect ratio and
    centering it. Falls back to a full-bleed rect when the image
    dimensions can't be probed.
    """
    try:
        from PIL import Image
        with Image.open(image_path) as im:
            img_w, img_h = im.size
    except Exception:
        return (0, 0, slide_w_emu, slide_h_emu)
    if img_w <= 0 or img_h <= 0:
        return (0, 0, slide_w_emu, slide_h_emu)
    slide_aspect = slide_w_emu / slide_h_emu
    img_aspect = img_w / img_h
    if img_aspect > slide_aspect:
        # Wider than slide — fit to slide width.
        width = slide_w_emu
        height = int(slide_w_emu / img_aspect)
        left = 0
        top = (slide_h_emu - height) // 2
    else:
        # Taller than (or equal to) slide — fit to slide height.
        height = slide_h_emu
        width = int(slide_h_emu * img_aspect)
        top = 0
        left = (slide_w_emu - width) // 2
    return (left, top, width, height)


def _make_video_poster(video_path: Path) -> Optional[Path]:
    """Best-effort poster frame for an embedded video — grabs
    a single still ~1 second in via ffmpeg. Returns the path on
    success, or None when ffmpeg is missing / extraction fails.
    The pptx exporter shows a blank media tile rather than crash
    when None is returned."""
    import shutil as _shutil
    import subprocess as _sp
    if not _shutil.which("ffmpeg"):
        return None
    try:
        poster = video_path.with_suffix(
            video_path.suffix + ".poster.png")
        if poster.exists():
            return poster
        proc = _sp.run(
            ["ffmpeg", "-y", "-ss", "1",
             "-i", str(video_path),
             "-frames:v", "1",
             str(poster)],
            capture_output=True, text=True, timeout=30)
        if proc.returncode == 0 and poster.exists():
            return poster
    except Exception:
        return None
    return None


def build_deck_entries(
    scenes: List[Any],
    title_card_dir: Optional[Path] = None,
    chapter_title: str = "",
    default_image_seconds: float = 4.0,
) -> Tuple[List[Path], List[float], List[str]]:
    """Walk the scenes and return three parallel lists ready for
    ``stitcher.stitch_clips``: paths, per-clip durations, and a
    short summary line per entry (for the UI's "skipped" report).

    When ``title_card_dir`` is set, a title card PNG is generated
    for each scene and prepended to that scene's clip.

    Scenes whose favorite clip is missing, a placeholder, or zero
    bytes are skipped. The returned ``skipped`` list explains each
    skip for the UI.
    """
    paths: List[Path] = []
    durations: List[float] = []
    skipped: List[str] = []
    for idx, scene in enumerate(scenes, start=1):
        clip = scene.favorite_clip()
        label = scene.name or f"Scene {idx}"
        if clip is None or not clip.file_path:
            skipped.append(f"{label}: no favorite output")
            continue
        path = Path(clip.file_path)
        if not path.exists() or path.stat().st_size == 0:
            skipped.append(f"{label}: file missing or empty")
            continue
        if getattr(clip, "is_placeholder", False):
            skipped.append(f"{label}: placeholder only")
            continue
        # Optional title card.
        if title_card_dir is not None:
            title_card_dir.mkdir(parents=True, exist_ok=True)
            card_path = (
                title_card_dir
                / f"title_{idx:03d}_{scene.id}.png")
            overline = (
                f"{chapter_title} · Scene {idx}"
                if chapter_title else f"Scene {idx}")
            render_title_card(
                card_path,
                title=scene.name or f"Scene {idx}",
                subtitle=(scene.description or "")[:280],
                overline=overline)
            paths.append(card_path)
            durations.append(TITLE_CARD_DURATION_SECONDS)
        paths.append(path)
        # Honor each clip's stored duration. Image stills carry the
        # writer's chosen display time; video clips carry the
        # actual length (ffmpeg ignores -t on already-encoded
        # videos so passing it through is harmless).
        durations.append(
            float(clip.duration_seconds or default_image_seconds))
    return paths, durations, skipped
